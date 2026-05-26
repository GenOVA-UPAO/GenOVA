"""File-type-aware text extractors for the RAG pipeline.

Each parser receives a filesystem path and returns plain UTF-8 text (or raises
`ParserError` if the file type isn't supported or extraction fails). Audio and
image parsers reuse the existing Groq Whisper/vision helpers — they hit the
network and should be timed out by the caller.
"""
from __future__ import annotations

import logging
import os
from collections.abc import Callable

logger = logging.getLogger(__name__)


class ParserError(RuntimeError):
    pass


_EXT_MAP: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".mp3": "audio",
    ".wav": "audio",
    ".m4a": "audio",
    ".aac": "audio",
    ".ogg": "audio",
    ".webm": "audio",
    ".jpg": "image",
    ".jpeg": "image",
    ".png": "image",
    ".gif": "image",
    ".webp": "image",
}


def detect_kind(filename: str) -> str | None:
    ext = os.path.splitext(filename)[1].lower()
    return _EXT_MAP.get(ext)


def _parse_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ParserError("pypdf not installed") from exc
    text_parts: list[str] = []
    reader = PdfReader(path)
    for page in reader.pages:
        try:
            text_parts.append(page.extract_text() or "")
        except Exception as exc:
            logger.debug("Failed to extract text from PDF page: %s", exc)
            continue
    return "\n".join(p for p in text_parts if p.strip())


def _parse_docx(path: str) -> str:
    try:
        import docx  # type: ignore
    except ImportError as exc:
        raise ParserError("python-docx not installed") from exc
    document = docx.Document(path)
    return "\n".join(p.text for p in document.paragraphs if p.text.strip())


def _parse_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise ParserError("python-pptx not installed") from exc
    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                parts.append(text)
    return "\n".join(parts)


def _parse_audio(path: str) -> str:
    from agents.audio_helpers import transcribir_audio
    return transcribir_audio(path)


def _parse_image(path: str) -> str:
    import base64

    from agents.llm_router import generar_vision

    with open(path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("ascii")

    ext = os.path.splitext(path)[1].lower().lstrip(".") or "png"
    mime = "jpeg" if ext == "jpg" else ext
    messages = [
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": (
                        "Describe esta imagen en español con detalle educativo: "
                        "qué muestra, qué conceptos académicos representa, datos visibles. "
                        "Máximo 250 palabras."
                    ),
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/{mime};base64,{b64}"},
                },
            ],
        }
    ]
    return generar_vision(messages, max_tokens=512)


_PARSERS: dict[str, Callable[[str], str]] = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "audio": _parse_audio,
    "image": _parse_image,
}


def extract_text(path: str, filename: str | None = None) -> str:
    """Extract plain text from a file at `path`. `filename` defaults to the
    basename of `path` and is used only for extension detection — useful when
    the file lives at a hashed/generated path."""
    name = filename or os.path.basename(path)
    kind = detect_kind(name)
    if not kind:
        raise ParserError(f"Unsupported file type: {name}")
    parser = _PARSERS[kind]
    try:
        text = parser(path)
    except ParserError:
        raise
    except Exception as exc:
        logger.exception("Parser failed for %s (%s)", name, kind)
        raise ParserError(f"Failed to extract text from {name}: {exc}") from exc
    return text.strip()
